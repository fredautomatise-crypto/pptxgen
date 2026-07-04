"""
Génère un fichier .pptx à partir d'un CoursePlan (schema.py), en reproduisant
le template visuel observé dans les modules existants (fond marine pour les
slides titre/section, fond blanc pour le contenu, badges orange, cercles
icônes colorés, polices Cambria/Calibri).

Usage direct :
    python pptx_generator.py cache/plan.json output/module.pptx
"""
import sys
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import config
from schema import CoursePlan, Slide

C = {k: RGBColor.from_string(v) for k, v in config.COLORS.items()}
ICON_COLORS = config.ICON_PALETTE
FONT_H = config.FONTS["header"]
FONT_B = config.FONTS["body"]

SW, SH = config.SLIDE_WIDTH_IN, config.SLIDE_HEIGHT_IN
MARGIN = 0.6


# ---------------------------------------------------------------------------
# Helpers bas niveau
# ---------------------------------------------------------------------------

def new_slide(prs: Presentation, bg_hex: str):
    layout = prs.slide_layouts[6]  # layout vide
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(bg_hex)
    bg.line.fill.background()
    bg.shadow.inherit = False
    # envoyer le rectangle de fond tout en bas de la pile z-order
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return slide


def add_text(slide, text, left, top, width, height, size=14, color="1E293B",
             bold=False, italic=False, font=FONT_B, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.05, wrap=True, shrink=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = RGBColor.from_string(color) if isinstance(color, str) else color
    return box


def add_rounded_rect(slide, left, top, width, height, fill_hex=None, line_hex=None,
                      line_w=1.0, radius=0.08, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = RGBColor.from_string(line_hex)
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_icon_circle(slide, icon_key, left, top, diameter, color_hex, glyph_color="FFFFFF"):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top),
                                     Inches(diameter), Inches(diameter))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = config.ICON_MAP.get(icon_key, "●")
    r.font.size = Pt(diameter * 26)
    r.font.color.rgb = RGBColor.from_string(glyph_color)
    return circle


def add_badge(slide, text, left, top, fill_hex=config.COLORS["accent"], text_hex="FFFFFF",
              size=10, width=None):
    width = width or (0.18 * len(text) + 0.4)
    badge = add_rounded_rect(slide, left, top, width, 0.32, fill_hex=fill_hex, radius=0.5)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT_B
    r.font.color.rgb = RGBColor.from_string(text_hex)
    return badge


def add_footer(slide, text, dark_bg=False):
    color = "64748B" if not dark_bg else "64748B"
    add_text(slide, text, MARGIN, SH - 0.45, SW - 2 * MARGIN, 0.3, size=9,
              color=color, italic=True, align=PP_ALIGN.LEFT)


def add_decorative_circles(slide, dark_hex="1B3A5C"):
    """Motif de cercles concentriques décoratifs (coin haut droit), présent
    sur les slides titre/section du template source."""
    specs = [(9.6, -1.2, 4.2), (10.4, -0.2, 2.6), (11.0, 0.6, 1.4)]
    for left, top, size in specs:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor.from_string(dark_hex)
        c.fill.transparency = 0
        c.line.fill.background()
        c.shadow.inherit = False


def slide_title(slide, text, top=0.55, size=26, color="1E293B", width=SW - 2 * MARGIN):
    add_text(slide, text, MARGIN, top, width, 0.7, size=size, color=color, bold=True, font=FONT_H)


def _estimate_lines(text: str, font_size: int, width_in: float, chars_per_inch_at_1pt: float = 138.5) -> int:
    """Estimation grossière (mais suffisante) du nombre de lignes qu'un texte
    en gras Cambria occupera dans une boîte de largeur donnée, pour ajuster
    dynamiquement la taille de police et la position des éléments en dessous."""
    import textwrap
    chars_per_line = max(int(width_in * chars_per_inch_at_1pt / font_size), 1)
    return len(textwrap.wrap(text, chars_per_line)) or 1


def item_label_desc(item):
    """
    Les items du plan JSON ont 5 champs (icon/label/value/text/description)
    mais le LLM ne remplit pas toujours exactement ceux attendus par chaque
    type de slide. Ce helper construit un (label, description) robuste,
    quel que soit le champ que l'IA a effectivement rempli.
    """
    label = item.label or item.text or item.value or ""
    if item.label:
        desc = item.description or (item.text if item.text != label else "")
    else:
        desc = item.description
    return label, desc


# ---------------------------------------------------------------------------
# Rendu par type de slide
# ---------------------------------------------------------------------------

def render_title(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_dark"])
    add_decorative_circles(s)
    add_badge(s, f"Module {plan.module_number} de {plan.module_total}  ·  {plan.product_name or config.BRAND_NAME}",
               MARGIN, 0.6, fill_hex=config.COLORS["accent"])

    title_text = plan.module_title.upper()
    title_w = 9.8
    # Choisit la plus grande taille de police qui tient sur ~2 lignes, sinon
    # réduit progressivement pour éviter tout chevauchement avec le sous-titre.
    font_size = 24
    for candidate in (44, 38, 32, 28, 24):
        if _estimate_lines(title_text, candidate, title_w) <= 2:
            font_size = candidate
            break
    n_lines = _estimate_lines(title_text, font_size, title_w)
    line_height = font_size * 0.0245  # pouces — marge généreuse, le rendu réel (PowerPoint/LibreOffice)
    title_box_h = max(n_lines * line_height + 0.15, 0.75)  # dépasse souvent l'estimation "théorique"
    title_top = 1.3
    add_text(s, title_text, MARGIN, title_top, title_w, title_box_h, size=font_size, bold=True,
              color="FFFFFF", font=FONT_H)

    subtitle_top = title_top + title_box_h + 0.4
    add_text(s, slide_data.subtitle or plan.subtitle, MARGIN, subtitle_top, 8.5, 0.6, size=20,
              color="E2E8F0", font=FONT_H, italic=True)

    desc = slide_data.title if slide_data.title and slide_data.title != plan.module_title else ""
    desc_top = subtitle_top + 0.65
    add_text(s, desc, MARGIN, desc_top, 7.5, 1.0, size=14, color="CBD5E1", font=FONT_B, line_spacing=1.3)

    add_text(s, f"{config.BRAND_NAME}  ·  {config.TAGLINE}", MARGIN, SH - 0.9,
              8, 0.4, size=11, color="64748B", font=FONT_B)
    return s


def render_section(prs, plan: CoursePlan, slide_data: Slide, index_hint: str = ""):
    s = new_slide(prs, config.COLORS["bg_dark"])
    add_decorative_circles(s)
    add_badge(s, index_hint or "SECTION", 0.6, 3.15, fill_hex=config.COLORS["accent"])
    add_text(s, slide_data.title, MARGIN, 3.6, 10, 1.2, size=36, bold=True, color="FFFFFF", font=FONT_H)
    return s


def render_objectives(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_light"])
    slide_title(s, slide_data.title or "Ce que vous allez maîtriser dans ce module")

    # Colonne gauche : objectifs à puces icônes
    left_w = 7.6
    y = 1.6
    row_h = 0.85
    for i, item in enumerate(slide_data.items):
        color = ICON_COLORS[i % len(ICON_COLORS)]
        add_icon_circle(s, item.icon, MARGIN, y, 0.42, color)
        add_text(s, item.text or item.description, MARGIN + 0.6, y - 0.03, left_w - 0.6, row_h,
                  size=13, color="1E293B", line_spacing=1.15)
        y += row_h
        if y > SH - 0.9:
            break

    # Colonne droite : prérequis + durée
    rx = MARGIN + left_w + 0.3
    rw = SW - rx - MARGIN
    box = add_rounded_rect(s, rx, 1.6, rw, 3.0, fill_hex="F8FAFC", line_hex="E2E8F0")
    add_text(s, "PRÉREQUIS", rx + 0.25, 1.8, rw - 0.5, 0.3, size=12, bold=True, color="C2370A")
    prereq = [it.text or it.description for it in slide_data.items if it.label.lower() == "prerequis"]
    if not prereq:
        prereq = [slide_data.column_left_title] if slide_data.column_left_title else []
    prereq = slide_data.column_left_items or prereq
    ty = 2.2
    for p in prereq[:5]:
        add_text(s, f"–  {p}", rx + 0.25, ty, rw - 0.5, 0.4, size=11.5, color="1E293B", line_spacing=1.1)
        ty += 0.42

    add_rounded_rect(s, rx, 4.85, rw, 0.9, fill_hex="FFF7ED", line_hex="F5C99B")
    add_text(s, "DURÉE ESTIMÉE", rx + 0.25, 5.0, rw - 0.5, 0.25, size=11, bold=True, color="C2370A")
    add_text(s, slide_data.column_right_title or slide_data.subtitle or "20 – 30 minutes",
              rx + 0.25, 5.3, rw - 0.5, 0.4, size=13, bold=True, color="1E293B")

    add_footer(s, f"{config.BRAND_NAME}  ·  Module {plan.module_number} – {plan.module_title}")
    return s


def render_stats(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_light"])
    slide_title(s, slide_data.title)
    if slide_data.subtitle:
        add_text(s, slide_data.subtitle, MARGIN, 1.35, SW - 2 * MARGIN, 0.6, size=13,
                  color="64748B", italic=True, line_spacing=1.2)
    n = max(len(slide_data.items), 1)
    gap = 0.4
    col_w = (SW - 2 * MARGIN - gap * (n - 1)) / n
    x = MARGIN
    for i, item in enumerate(slide_data.items):
        color = ICON_COLORS[i % len(ICON_COLORS)]
        label, desc = item_label_desc(item)
        add_icon_circle(s, item.icon, x, 2.3, 0.55, color)
        add_text(s, item.value, x, 3.05, col_w, 0.7, size=30, bold=True, color=color, font=FONT_H)
        add_text(s, label, x, 3.75, col_w, 0.5, size=12, bold=True, color="1E293B", line_spacing=1.1)
        add_text(s, desc, x, 4.25, col_w, 1.2, size=10.5, color="64748B", line_spacing=1.15)
        x += col_w + gap
    add_footer(s, f"{config.BRAND_NAME}  ·  Module {plan.module_number} – {plan.module_title}")
    return s


def render_icon_list(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_light"])
    slide_title(s, slide_data.title)
    y = 1.55
    row_h = min(1.05, (SH - 1.9) / max(len(slide_data.items), 1))
    for i, item in enumerate(slide_data.items):
        color = ICON_COLORS[i % len(ICON_COLORS)]
        label, desc = item_label_desc(item)
        add_icon_circle(s, item.icon, MARGIN, y + 0.05, 0.45, color)
        add_text(s, label, MARGIN + 0.65, y, SW - 2 * MARGIN - 0.65, 0.35, size=14,
                  bold=True, color="1E293B")
        add_text(s, desc, MARGIN + 0.65, y + 0.36, SW - 2 * MARGIN - 0.65, 0.5,
                  size=11.5, color="64748B", line_spacing=1.1)
        y += row_h
    add_footer(s, f"{config.BRAND_NAME}  ·  Module {plan.module_number} – {plan.module_title}")
    return s


def render_icon_grid(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_light"])
    slide_title(s, slide_data.title)
    items = slide_data.items
    cols = 3 if len(items) > 4 else 2
    rows = (len(items) + cols - 1) // cols
    gap = 0.3
    card_w = (SW - 2 * MARGIN - gap * (cols - 1)) / cols
    card_h = min(1.9, (SH - 1.9 - gap * (rows - 1)) / max(rows, 1))
    for i, item in enumerate(items):
        r, c = divmod(i, cols)
        x = MARGIN + c * (card_w + gap)
        y = 1.55 + r * (card_h + gap)
        color = ICON_COLORS[i % len(ICON_COLORS)]
        label, desc = item_label_desc(item)
        add_rounded_rect(s, x, y, card_w, card_h, fill_hex="F8FAFC", line_hex="E2E8F0")
        add_icon_circle(s, item.icon, x + 0.2, y + 0.2, 0.4, color)
        add_text(s, label, x + 0.2, y + 0.75, card_w - 0.4, 0.35, size=12.5, bold=True, color="1E293B")
        add_text(s, desc, x + 0.2, y + 1.1, card_w - 0.4, card_h - 1.2, size=10,
                  color="64748B", line_spacing=1.1)
    add_footer(s, f"{config.BRAND_NAME}  ·  Module {plan.module_number} – {plan.module_title}")
    return s


def render_two_column(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_light"])
    slide_title(s, slide_data.title)
    col_w = (SW - 2 * MARGIN - 0.4) / 2
    cols = [
        (MARGIN, slide_data.column_left_title, slide_data.column_left_items, "990011", "FCF6F5"),
        (MARGIN + col_w + 0.4, slide_data.column_right_title, slide_data.column_right_items, "0F766E", "F0FDFA"),
    ]
    for x, title, bullets, accent, bg in cols:
        add_rounded_rect(s, x, 1.55, col_w, SH - 2.3, fill_hex=bg, line_hex="E2E8F0")
        add_text(s, title, x + 0.25, 1.75, col_w - 0.5, 0.4, size=14, bold=True, color=accent)
        y = 2.25
        for b in bullets:
            add_text(s, f"–  {b}", x + 0.25, y, col_w - 0.5, 0.5, size=11.5, color="1E293B", line_spacing=1.15)
            y += 0.5
    add_footer(s, f"{config.BRAND_NAME}  ·  Module {plan.module_number} – {plan.module_title}")
    return s


def render_table(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_light"])
    slide_title(s, slide_data.title)
    items = slide_data.items
    n = max(len(items), 1)
    rows, cols = n + 1, 2
    table_shape = s.shapes.add_table(rows, cols, Inches(MARGIN), Inches(1.6),
                                      Inches(SW - 2 * MARGIN), Inches(min(4.8, 0.5 * rows)))
    table = table_shape.table
    table.columns[0].width = Inches((SW - 2 * MARGIN) * 0.35)
    table.columns[1].width = Inches((SW - 2 * MARGIN) * 0.65)
    headers = [slide_data.column_left_title or "Élément", slide_data.column_right_title or "Description"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(config.COLORS["bg_dark"])
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("FFFFFF")
        run.font.size = Pt(12)
        run.font.name = FONT_B
    for r, item in enumerate(items, start=1):
        for c, val in enumerate([item.label, item.description or item.text]):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("F8FAFC" if r % 2 else "FFFFFF")
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(11)
            run.font.name = FONT_B
            run.font.color.rgb = RGBColor.from_string("1E293B")
    add_footer(s, f"{config.BRAND_NAME}  ·  Module {plan.module_number} – {plan.module_title}")
    return s


def render_process_flow(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_light"])
    slide_title(s, slide_data.title)
    items = slide_data.items
    n = max(len(items), 1)
    gap_arrow = 0.4
    box_w = (SW - 2 * MARGIN - gap_arrow * (n - 1)) / n - 0.1
    y = 3.0
    x = MARGIN
    for i, item in enumerate(items):
        color = ICON_COLORS[i % len(ICON_COLORS)]
        label, desc = item_label_desc(item)
        add_rounded_rect(s, x, y, box_w, 1.6, fill_hex="F8FAFC", line_hex=color, line_w=1.5)
        add_icon_circle(s, item.icon, x + box_w / 2 - 0.22, y + 0.18, 0.44, color)
        add_text(s, label, x, y + 0.75, box_w, 0.35, size=11.5, bold=True, color="1E293B",
                  align=PP_ALIGN.CENTER)
        add_text(s, desc, x, y + 1.1, box_w, 0.5, size=9.5, color="64748B",
                  align=PP_ALIGN.CENTER, line_spacing=1.05)
        x += box_w
        if i < n - 1:
            add_text(s, "→", x, y + 0.5, gap_arrow, 0.6, size=22, bold=True,
                      color=config.COLORS["accent"], align=PP_ALIGN.CENTER)
            x += gap_arrow
    add_footer(s, f"{config.BRAND_NAME}  ·  Module {plan.module_number} – {plan.module_title}")
    return s


def render_conclusion(prs, plan: CoursePlan, slide_data: Slide):
    s = new_slide(prs, config.COLORS["bg_dark"])
    add_decorative_circles(s)
    add_badge(s, "RÉCAPITULATIF", MARGIN, 0.6, fill_hex=config.COLORS["accent"])
    add_text(s, slide_data.title or "Ce qu'il faut retenir", MARGIN, 1.15, 10, 0.8, size=32,
              bold=True, color="FFFFFF", font=FONT_H)
    y = 2.3
    for item in slide_data.items:
        add_icon_circle(s, "check", MARGIN, y, 0.35, config.COLORS["accent"])
        add_text(s, item.text or item.description, MARGIN + 0.55, y - 0.02, 10.5, 0.5,
                  size=13, color="E2E8F0", line_spacing=1.15)
        y += 0.55
    add_text(s, f"{config.BRAND_NAME}  ·  {config.TAGLINE}", MARGIN, SH - 0.7,
              8, 0.4, size=11, color="64748B")
    return s


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "objectives": render_objectives,
    "stats": render_stats,
    "icon_list": render_icon_list,
    "icon_grid": render_icon_grid,
    "two_column": render_two_column,
    "table": render_table,
    "process_flow": render_process_flow,
    "conclusion": render_conclusion,
}


def generate_pptx(plan: CoursePlan, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    section_count = 0
    for slide_data in plan.slides:
        renderer = RENDERERS.get(slide_data.type)
        if renderer is None:
            print(f"[warn] type de slide inconnu ignoré: {slide_data.type}")
            continue
        if slide_data.type == "section":
            section_count += 1
            rendered_slide = renderer(prs, plan, slide_data, index_hint=f"SECTION {section_count}")
        else:
            rendered_slide = renderer(prs, plan, slide_data)

        # Le champ `speaker_notes` (texte détaillé généré par le LLM, voir
        # schema.py) est injecté comme note de présentateur PowerPoint —
        # visible en mode "Présentateur" ou à l'impression, mais jamais
        # affiché à l'écran pendant la présentation. C'est l'équivalent
        # "sans audio" du script de narration : Fred garde un support texte
        # complet pour présenter/enregistrer lui-même s'il le souhaite.
        if slide_data.speaker_notes and rendered_slide is not None:
            rendered_slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    import content_planner
    if len(sys.argv) < 3:
        print("Usage: python pptx_generator.py plan.json output.pptx")
        sys.exit(1)
    plan_path, out_path = sys.argv[1], sys.argv[2]
    plan = content_planner.load_plan(plan_path)
    path = generate_pptx(plan, out_path)
    print(f"PPTX généré : {path}")
